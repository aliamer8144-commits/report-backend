from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from typing import Optional
from io import BytesIO
from pathlib import Path
import os
import re
import time
import uuid
import threading
from dotenv import load_dotenv

# Load environment variables from .env.local
load_dotenv(dotenv_path=".env.local")

from pptx import Presentation  # python-pptx
from urllib.parse import quote
import requests

# --- Aspose Token Cache ---
_cached_token: str | None = None
_token_expiry: float = 0
_token_lock = threading.Lock()

# --- Font upload state (fonts persist in Aspose Cloud Storage forever) ---
_fonts_initialized = False
_fonts_lock = threading.Lock()


class ReportPayload(BaseModel):
    SERVICE_CODE: str
    ID_NUMBER: str
    NAME_AR: str
    NAME_EN: str
    DAYS_COUNT: int
    ENTRY_DATE_GREGORIAN: str
    EXIT_DATE_GREGORIAN: str
    ENTRY_DATE_HIJRI: Optional[str] = None
    EXIT_DATE_HIJRI: Optional[str] = None
    REPORT_ISSUE_DATE: str
    NATIONALITY_AR: str
    NATIONALITY_EN: str
    DOCTOR_NAME_AR: str
    DOCTOR_NAME_EN: str
    JOB_TITLE_AR: str
    JOB_TITLE_EN: str
    HOSPITAL_NAME_AR: str
    HOSPITAL_NAME_EN: str
    PRINT_DATE: str
    PRINT_TIME: str


def get_template_path() -> Path:
    # Allow override via env for both path and name
    env_path = os.getenv("PPTX_TEMPLATE_PATH")
    if env_path:
        p = Path(env_path)
        if p.exists():
            return p

    template_name = os.getenv("PPTX_TEMPLATE_NAME", "report_template.pptx")
    current_dir = Path(__file__).resolve().parent
    
    # Priority candidates based on your project structure
    candidates = [
        # Local development (CWD is backend/)
        Path("public/templates") / template_name,
        # Vercel deployment paths
        Path("/var/task/backend/public/templates") / template_name,
        Path("/var/task/public/templates") / template_name,
        # Relative to current script
        current_dir / "public" / "templates" / template_name,
        current_dir.parent / "public" / "templates" / template_name,
        # Common Vercel folder structure
        Path("/var/task/api/public/templates") / template_name,
        Path("/var/task/backend/api/public/templates") / template_name,
    ]
    
    for candidate in candidates:
        if candidate.exists():
            return candidate
            
    # Search recursively as a fallback
    try:
        for match in Path.cwd().rglob(template_name):
            return match
    except Exception:
        pass
        
    return candidates[0]


def format_date_dd_mm_yyyy(value: Optional[str]) -> Optional[str]:
    if value is None:
        return None
    s = str(value).strip()
    # Find first occurrence of YYYY-MM-DD or YYYY/MM/DD anywhere in the string (e.g., ISO timestamps)
    m = re.search(r"(\d{4})[-/](\d{1,2})[-/](\d{1,2})", s)
    if not m:
        return s
    yyyy, mm, dd = m.groups()
    mm = mm.zfill(2)
    dd = dd.zfill(2)
    return f"{dd}-{mm}-{yyyy}"


def load_template_presentation() -> Presentation:
    """Load template as Presentation from local file or URL."""
    # Priority 1: Local file
    local_path = get_template_path()
    if local_path and local_path.exists():
        try:
            return Presentation(str(local_path))
        except Exception as e:
            print(f"Error loading local template {local_path}: {e}")

    # Priority 2: External URL fallback
    template_url = os.getenv("PPTX_TEMPLATE_URL")
    if not template_url:
        raise HTTPException(
            status_code=500, 
            detail="Local template not found and PPTX_TEMPLATE_URL is not set"
        )
    
    try:
        resp = requests.get(template_url, timeout=20)
        if resp.status_code != 200:
            raise HTTPException(status_code=500, detail=f"Failed to fetch template from URL: {resp.status_code}")
        return Presentation(BytesIO(resp.content))
    except requests.RequestException as e:
        raise HTTPException(status_code=500, detail=f"Error fetching template from URL: {str(e)}")


def replace_placeholders(prs: Presentation, mapping: dict):
    # Replace text placeholders in all shapes across all slides
    # Do replacements per-run to preserve formatting (font color/size)
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if hasattr(shape, "text_frame") and shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text = run.text or ""
                            new_text = text
                            for key, value in mapping.items():
                                new_text = new_text.replace(f"{{{{{key}}}}}", str(value) if value is not None else "")
                            if new_text != text:
                                run.text = new_text  # preserves run formatting
            except Exception:
                # Skip shapes that fail to process to avoid taking down the request
                continue


app = FastAPI(title="PPTX Generator Service")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # adjust in production
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


@app.get("/")
@app.get("/health")
def health():
    template_url = os.getenv("PPTX_TEMPLATE_URL")
    template_name = os.getenv("PPTX_TEMPLATE_NAME", "report_template.pptx")
    p = get_template_path()
    return {
        "status": "ok",
        "template_name": template_name,
        "template_found": p.exists(),
        "resolved_path": str(p),
        "template_url_configured": bool(template_url)
    }


@app.get("/debug-template")
def debug_template():
    p = get_template_path()
    candidates = []
    try:
        roots = [Path("/var/task"), Path.cwd()]
        found = []
        for root in roots:
            try:
                for match in root.rglob("report_template.pptx"):
                    found.append(str(match))
            except Exception:
                continue
    except Exception:
        found = []
    checks = {
        "/var/task/backend/public/templates": (Path("/var/task/backend/public/templates").exists()),
        "/var/task/public/templates": (Path("/var/task/public/templates").exists()),
    }
    template_url = os.getenv("PPTX_TEMPLATE_URL")
    return {
        "resolved_path": str(p),
        "exists": p.exists(),
        "cwd": str(Path.cwd()),
        "file_dir": str(Path(__file__).resolve().parent),
        "found_candidates": found,
        "dir_checks": checks,
        "template_url": template_url or None,
    }


@app.post("/generate-pptx")
def generate_pptx(payload: ReportPayload):
    try:
        prs = load_template_presentation()

        # Build mapping from placeholders to values.
        mapping = {
            "SERVICE_CODE": payload.SERVICE_CODE,
            "ID_NUMBER": payload.ID_NUMBER,
            "NAME_AR": payload.NAME_AR,
            "NAME_EN": payload.NAME_EN,
            "DAYS_COUNT": payload.DAYS_COUNT,
            "ENTRY_DATE_GREGORIAN": format_date_dd_mm_yyyy(payload.ENTRY_DATE_GREGORIAN),
            "EXIT_DATE_GREGORIAN": format_date_dd_mm_yyyy(payload.EXIT_DATE_GREGORIAN),
            "ENTRY_DATE_HIJRI": format_date_dd_mm_yyyy(payload.ENTRY_DATE_HIJRI),
            "EXIT_DATE_HIJRI": format_date_dd_mm_yyyy(payload.EXIT_DATE_HIJRI),
            "REPORT_ISSUE_DATE": format_date_dd_mm_yyyy(payload.REPORT_ISSUE_DATE),
            "NATIONALITY_AR": payload.NATIONALITY_AR,
            "NATIONALITY_EN": payload.NATIONALITY_EN,
            "DOCTOR_NAME_AR": payload.DOCTOR_NAME_AR,
            "DOCTOR_NAME_EN": payload.DOCTOR_NAME_EN,
            "JOB_TITLE_AR": payload.JOB_TITLE_AR,
            "JOB_TITLE_EN": payload.JOB_TITLE_EN,
            "HOSPITAL_NAME_AR": payload.HOSPITAL_NAME_AR,
            "HOSPITAL_NAME_EN": payload.HOSPITAL_NAME_EN,
            "PRINT_DATE": format_date_dd_mm_yyyy(payload.PRINT_DATE),
            "PRINT_TIME": payload.PRINT_TIME,
        }

        replace_placeholders(prs, mapping)

        buf = BytesIO()
        prs.save(buf)
        buf.seek(0)

        filename = "sickLeaves.pptx"
        # HTTP headers must be latin-1 encodable in Starlette; use RFC5987 filename*
        ascii_fallback = "sickLeaves.pptx"
        content_disposition = (
            f"attachment; filename=\"{ascii_fallback}\"; filename*=UTF-8''{quote(filename)}"
        )
        return StreamingResponse(
            buf,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": content_disposition
            },
        )
    except Exception as e:
        # Log server-side for debugging
        import traceback
        print("[generate-pptx] Error:", e)
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


ASPOSE_TOKEN_URL = "https://api.aspose.cloud/connect/token"
ASPOSE_SLIDES_API = "https://api.aspose.cloud/v3.0/slides"

# Font files directory - uploaded to Aspose Cloud Storage for PPTX→PDF conversion
# These are full-coverage Arabic fonts (renamed from Google Fonts) that replace the
# template's embedded font subsets. Aspose automatically matches fonts by internal name.
# Dubai → renamed Tajawal, Tahoma → renamed Noto Sans Arabic (255 Arabic glyphs)
_FONTS_DIR = Path(__file__).resolve().parent / "fonts"
FONT_STORAGE_FOLDER = "fonts"
REQUIRED_FONTS = [
    "Dubai-Regular.ttf",
    "Dubai-Bold.ttf",
    "Tahoma-Regular.ttf",
    "Tahoma-Bold.ttf",
]


def get_aspose_token() -> str:
    """Get JWT access token from Aspose Cloud (cached for ~58 minutes)."""
    global _cached_token, _token_expiry
    
    with _token_lock:
        if _cached_token and time.time() < _token_expiry:
            return _cached_token
    
    client_id = os.getenv("ASPOSE_APP_SID")
    client_secret = os.getenv("ASPOSE_APP_KEY")
    if not client_id or not client_secret:
        raise HTTPException(status_code=500, detail="ASPOSE_APP_SID and ASPOSE_APP_KEY must be configured")
    
    resp = requests.post(
        ASPOSE_TOKEN_URL,
        data={
            "grant_type": "client_credentials",
            "client_id": client_id,
            "client_secret": client_secret,
        },
        timeout=30,
    )
    if resp.status_code != 200:
        raise HTTPException(status_code=500, detail=f"Failed to get Aspose token: {resp.text[:300]}")
    
    token = resp.json()["access_token"]
    with _token_lock:
        _cached_token = token
        _token_expiry = time.time() + 3500  # ~58 min (tokens valid for 1 hour)
    return token


def _ensure_fonts_once(token: str):
    """Upload fonts to Aspose Cloud Storage exactly once per server instance.
    
    Fonts persist in storage permanently. This runs once on first PDF request.
    No API calls on subsequent requests (in-memory flag).
    """
    global _fonts_initialized
    if _fonts_initialized:
        return
    
    with _fonts_lock:
        # Double-check after acquiring lock
        if _fonts_initialized:
            return
        
        headers = {"Authorization": f"Bearer {token}"}
        for font_file in REQUIRED_FONTS:
            font_path = _FONTS_DIR / font_file
            if not font_path.exists():
                print(f"[fonts] Font file not found: {font_path}")
                continue
            try:
                with open(font_path, "rb") as f:
                    font_bytes = f.read()
                storage_path = f"{FONT_STORAGE_FOLDER}/{font_file}"
                resp = requests.put(
                    f"{ASPOSE_SLIDES_API}/storage/file/{storage_path}",
                    headers=headers,
                    data=font_bytes,
                    timeout=30,
                )
                if resp.status_code in (200, 201):
                    print(f"[fonts] Uploaded '{font_file}'")
                else:
                    print(f"[fonts] Upload '{font_file}': {resp.status_code}")
            except Exception as e:
                print(f"[fonts] Error uploading '{font_file}': {e}")
        
        _fonts_initialized = True


def convert_pptx_to_pdf_with_aspose(pptx_bytes: bytes) -> bytes:
    """Convert PPTX bytes to PDF using Aspose.Slides Cloud API.
    
    Optimized to use minimum API calls:
    - Token: cached in memory (0 extra calls after first)
    - Fonts: uploaded once ever (0 calls on subsequent requests)
    - Upload PPTX: 1 PUT call
    - Convert to PDF: 1 POST call
    - Delete: skipped (file overwritten next time)
    
    Total per PDF: 2 API calls (upload + convert)
    """
    token = get_aspose_token()
    auth_header = {"Authorization": f"Bearer {token}"}
    
    # Upload fonts once (no-op on subsequent requests)
    _ensure_fonts_once(token)

    # Generate unique filename to avoid race conditions
    file_id = uuid.uuid4().hex[:8]
    pptx_filename = f"r_{file_id}.pptx"

    # Step 1: Upload PPTX to Aspose Cloud Storage
    upload_resp = requests.put(
        f"{ASPOSE_SLIDES_API}/storage/file/{pptx_filename}",
        headers=auth_header,
        data=pptx_bytes,
        timeout=60,
    )
    if upload_resp.status_code not in (200, 201):
        raise HTTPException(
            status_code=500,
            detail=f"Failed to upload PPTX: {upload_resp.status_code} - {upload_resp.text[:300]}"
        )

    # Step 2: Convert PPTX to PDF with high quality options
    json_headers = {**auth_header, "Content-Type": "application/json"}
    pdf_options = {
        "format": "pdf",
        "options": {
            "jpegQuality": 100,
            "sufficientResolution": 300,
            "embedFullFonts": True,
            "saveMetafilesAsPng": True,
            "drawSlidesFrame": False,
            "fontFolders": [FONT_STORAGE_FOLDER],
        }
    }

    convert_resp = requests.post(
        f"{ASPOSE_SLIDES_API}/{pptx_filename}/pdf?withOptions=true",
        headers=json_headers,
        json=pdf_options,
        timeout=120,
    )
    
    # Fire-and-forget: delete uploaded file in background to save response time
    try:
        threading.Thread(
            target=lambda: requests.delete(
                f"{ASPOSE_SLIDES_API}/storage/file/{pptx_filename}",
                headers=auth_header,
                timeout=15,
            ),
            daemon=True,
        ).start()
    except Exception:
        pass

    if convert_resp.status_code != 200:
        raise HTTPException(
            status_code=500,
            detail=f"Failed to convert PPTX to PDF: {convert_resp.status_code} - {convert_resp.text[:300]}"
        )

    return convert_resp.content


@app.post("/generate-pdf")
def generate_pdf(payload: ReportPayload):
    try:
        prs = load_template_presentation()

        mapping = {
            "SERVICE_CODE": payload.SERVICE_CODE,
            "ID_NUMBER": payload.ID_NUMBER,
            "NAME_AR": payload.NAME_AR,
            "NAME_EN": payload.NAME_EN,
            "DAYS_COUNT": payload.DAYS_COUNT,
            "ENTRY_DATE_GREGORIAN": format_date_dd_mm_yyyy(payload.ENTRY_DATE_GREGORIAN),
            "EXIT_DATE_GREGORIAN": format_date_dd_mm_yyyy(payload.EXIT_DATE_GREGORIAN),
            "ENTRY_DATE_HIJRI": format_date_dd_mm_yyyy(payload.ENTRY_DATE_HIJRI),
            "EXIT_DATE_HIJRI": format_date_dd_mm_yyyy(payload.EXIT_DATE_HIJRI),
            "REPORT_ISSUE_DATE": format_date_dd_mm_yyyy(payload.REPORT_ISSUE_DATE),
            "NATIONALITY_AR": payload.NATIONALITY_AR,
            "NATIONALITY_EN": payload.NATIONALITY_EN,
            "DOCTOR_NAME_AR": payload.DOCTOR_NAME_AR,
            "DOCTOR_NAME_EN": payload.DOCTOR_NAME_EN,
            "JOB_TITLE_AR": payload.JOB_TITLE_AR,
            "JOB_TITLE_EN": payload.JOB_TITLE_EN,
            "HOSPITAL_NAME_AR": payload.HOSPITAL_NAME_AR,
            "HOSPITAL_NAME_EN": payload.HOSPITAL_NAME_EN,
            "PRINT_DATE": format_date_dd_mm_yyyy(payload.PRINT_DATE),
            "PRINT_TIME": payload.PRINT_TIME,
        }

        replace_placeholders(prs, mapping)

        # Save filled PPTX to in-memory buffer
        pptx_buf = BytesIO()
        prs.save(pptx_buf)
        pptx_buf.seek(0)

        # Convert PPTX -> PDF using Aspose.Slides Cloud API (high quality)
        pdf_bytes = convert_pptx_to_pdf_with_aspose(pptx_buf.getvalue())

        filename = "sickLeaves.pdf"
        ascii_fallback = "sickLeaves.pdf"
        cd = f"attachment; filename=\"{ascii_fallback}\"; filename*=UTF-8''{quote(filename)}"
        return StreamingResponse(
            BytesIO(pdf_bytes),
            media_type="application/pdf",
            headers={"Content-Disposition": cd},
        )
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        print("[generate-pdf] Error:", e)
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))
